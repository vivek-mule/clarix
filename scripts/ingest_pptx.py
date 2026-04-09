"""
scripts/ingest_pptx.py — One-time offline ingestion of .pptx course content into Pinecone.

Usage:
    # Single file
    python scripts/ingest_pptx.py --path slides/physics_ch1.pptx --subject physics

    # Entire folder (processes every .pptx inside)
    python scripts/ingest_pptx.py --path slides/ --subject physics

Behaviour:
    - Extracts text from every slide (title, body, bullets, speaker notes).
    - Each slide = 1 chunk.
    - Embeds locally with all-MiniLM-L6-v2 (384-dim, no API).
    - Upserts into Pinecone under a namespace matching the subject name.
    - Deterministic IDs ("filename_slide3") make re-runs idempotent.
    - Slides with no extractable text are skipped with a warning.
"""

from __future__ import annotations

import argparse
import os
import re
import sys
from pathlib import Path

# ── Ensure the backend package is importable ────────────────
# When run from the project root:  python scripts/ingest_pptx.py ...
BACKEND_DIR = Path(__file__).resolve().parent.parent / "backend"
sys.path.insert(0, str(BACKEND_DIR))

from pptx import Presentation  # python-pptx

from app.core.config import settings
from app.core.embeddings import embed_texts
from app.core.pinecone_client import index

# ── Constants ───────────────────────────────────────────────
UPSERT_BATCH_SIZE = 50  # Pinecone recommends batches ≤ 100


# ── Helpers ─────────────────────────────────────────────────

def _sanitize_id(text: str) -> str:
    """Turn a filename into a safe Pinecone vector ID component."""
    stem = Path(text).stem                        # remove extension
    slug = re.sub(r"[^a-zA-Z0-9]+", "_", stem)    # keep alphanumeric + _
    return slug.strip("_").lower()


def _infer_difficulty(slide_index: int, total_slides: int) -> str:
    """Map slide position to a rough difficulty tier."""
    ratio = slide_index / max(total_slides, 1)
    if ratio < 0.33:
        return "beginner"
    elif ratio < 0.66:
        return "intermediate"
    else:
        return "advanced"


def _extract_slide_text(slide) -> tuple[str, str]:
    """
    Return (title, body_text) from a single slide.
    body_text includes all shape text + speaker notes.
    """
    title = ""
    body_parts: list[str] = []

    for shape in slide.shapes:
        # Title placeholder
        if shape.has_text_frame:
            if shape.shape_id == slide.shapes.title and slide.shapes.title is not None:
                title = shape.text_frame.text.strip()
            else:
                text = shape.text_frame.text.strip()
                if text:
                    body_parts.append(text)

        # Tables
        if shape.has_table:
            for row in shape.table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.replace("|", "").strip():
                    body_parts.append(row_text)

    # Slide title fallback: use the dedicated title placeholder
    if not title and slide.shapes.title is not None:
        title = slide.shapes.title.text.strip()

    # Speaker notes
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            body_parts.append(f"[Speaker Notes] {notes}")

    body_text = "\n".join(body_parts)
    return title, body_text


# ── Per-file processing ────────────────────────────────────

def process_pptx(filepath: Path, subject: str) -> int:
    """
    Process a single .pptx file and upsert its slides to Pinecone.
    Returns the number of vectors upserted.
    """
    filename = filepath.name
    file_slug = _sanitize_id(filename)
    namespace = subject.strip().lower()

    print(f"\n{'━' * 60}")
    print(f"📂  Processing: {filename}")
    print(f"    Namespace:  {namespace}")
    print(f"{'━' * 60}")

    prs = Presentation(str(filepath))
    total_slides = len(prs.slides)

    vectors_to_upsert: list[dict] = []
    skipped = 0

    for idx, slide in enumerate(prs.slides, start=1):
        slide_title, slide_body = _extract_slide_text(slide)
        combined_text = f"{slide_title}\n{slide_body}".strip()

        if not combined_text:
            print(f"    ⚠  Slide {idx}/{total_slides} — no text, skipping")
            skipped += 1
            continue

        vector_id = f"{file_slug}_slide{idx}"
        difficulty = _infer_difficulty(idx, total_slides)

        metadata = {
            "source":     filename,
            "slide":      idx,
            "topic":      slide_title or f"Slide {idx}",
            "difficulty": difficulty,
            "subject":    namespace,
            "text":       combined_text[:40_000],  # Pinecone metadata limit safety
        }

        vectors_to_upsert.append({
            "id":       vector_id,
            "text":     combined_text,
            "metadata": metadata,
        })

        print(f"    ✔  Slide {idx}/{total_slides}  id={vector_id}  "
              f"difficulty={difficulty}  chars={len(combined_text)}")

    if not vectors_to_upsert:
        print(f"    ⛔  No vectors to upsert from {filename}")
        return 0

    # ── Embed all slide texts in one batch ──────────────────
    print(f"\n    🧠  Embedding {len(vectors_to_upsert)} slides locally …")
    texts = [v["text"] for v in vectors_to_upsert]
    embeddings = embed_texts(texts)

    # ── Upsert to Pinecone in batches ───────────────────────
    pinecone_vectors = [
        {
            "id":       v["id"],
            "values":   emb,
            "metadata": v["metadata"],
        }
        for v, emb in zip(vectors_to_upsert, embeddings)
    ]

    upserted = 0
    for i in range(0, len(pinecone_vectors), UPSERT_BATCH_SIZE):
        batch = pinecone_vectors[i : i + UPSERT_BATCH_SIZE]
        index.upsert(vectors=batch, namespace=namespace)
        upserted += len(batch)
        print(f"    📤  Upserted batch {i // UPSERT_BATCH_SIZE + 1}  "
              f"({len(batch)} vectors)")

    print(f"\n    ✅  Done — {upserted} upserted, {skipped} skipped")
    return upserted


# ── CLI entry point ─────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Ingest .pptx course slides into Pinecone for the Adaptive Learning Platform.",
    )
    parser.add_argument(
        "--path",
        required=True,
        help="Path to a single .pptx file OR a folder containing .pptx files.",
    )
    parser.add_argument(
        "--subject",
        required=True,
        help="Subject / namespace to store vectors under (e.g. 'physics', 'math').",
    )
    args = parser.parse_args()

    target = Path(args.path).resolve()

    # ── Collect .pptx files ─────────────────────────────────
    if target.is_file():
        if target.suffix.lower() != ".pptx":
            print(f"❌  {target} is not a .pptx file")
            sys.exit(1)
        pptx_files = [target]
    elif target.is_dir():
        pptx_files = sorted(target.glob("*.pptx"))
        if not pptx_files:
            print(f"❌  No .pptx files found in {target}")
            sys.exit(1)
    else:
        print(f"❌  Path does not exist: {target}")
        sys.exit(1)

    print(f"🚀  Found {len(pptx_files)} .pptx file(s) to process")
    print(f"    Pinecone index: {settings.pinecone_index_name}")
    print(f"    Subject/namespace: {args.subject}")

    total_upserted = 0
    for pptx_path in pptx_files:
        try:
            total_upserted += process_pptx(pptx_path, args.subject)
        except Exception as e:
            print(f"\n    ❌  Error processing {pptx_path.name}: {e}")
            continue

    print(f"\n{'━' * 60}")
    print(f"🏁  All done — {total_upserted} total vectors upserted "
          f"across {len(pptx_files)} file(s)")
    print(f"{'━' * 60}")


if __name__ == "__main__":
    main()
